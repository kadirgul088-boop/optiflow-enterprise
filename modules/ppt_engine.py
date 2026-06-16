from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
SLATE = RGBColor(51, 65, 85)
GRAY = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
LINE = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(220, 38, 38)
AMBER = RGBColor(217, 119, 6)
BLACK = RGBColor(17, 24, 39)


def _money(value):
    try:
        return f"{float(value):,.0f} TL".replace(",", ".")
    except Exception:
        return f"{value} TL"


def _pct(value):
    try:
        return f"%{float(value):.1f}"
    except Exception:
        return f"%{value}"


def _safe(data, key, default=0):
    try:
        return data.get(key, default)
    except Exception:
        return default


def _float(value, default=0):
    try:
        return float(value)
    except Exception:
        return default


def _color_for(label, value):
    value = _float(value)

    if label == "wait":
        return RED if value > 20 else AMBER if value > 12 else GREEN

    if label == "oee":
        return GREEN if value >= 80 else AMBER if value >= 65 else RED

    if label == "defect":
        return GREEN if value <= 3 else AMBER if value <= 7 else RED

    if label == "balance":
        return RED if value > 15 else AMBER if value > 8 else GREEN

    if label == "roi":
        return GREEN if value >= 150 else AMBER if value >= 75 else RED

    return BLUE


def _status_text(label, value):
    color = _color_for(label, value)
    if color == GREEN:
        return "GÜÇLÜ"
    if color == AMBER:
        return "DİKKAT"
    return "KRİTİK"


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _shape(slide, shape_type, x, y, w, h, fill=None, line=None, radius=True):
    shp = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill

    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line

    return shp


def _textbox(slide, text, x, y, w, h, size=12, color=BLACK, bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT

    return box


def _footer(slide, page_no):
    _textbox(slide, "OPTIFLOW ENTERPRISE | CONFIDENTIAL", 0.55, 7.12, 5.0, 0.18, 7.5, GRAY)
    _textbox(slide, f"{page_no}", 12.4, 7.12, 0.4, 0.18, 7.5, GRAY, align="right")


def _header(slide, title, kicker=None, page_no=None):
    _textbox(slide, title, 0.55, 0.35, 8.8, 0.35, 21, NAVY, True)
    if kicker:
        _textbox(slide, kicker, 0.57, 0.78, 9.6, 0.25, 9.5, GRAY)

    _shape(slide, MSO_SHAPE.RECTANGLE, 0.55, 1.13, 12.25, 0.02, LINE)

    if page_no is not None:
        _footer(slide, page_no)


def _background(slide):
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, LIGHT)


def _metric_card(slide, title, value, subtitle, x, y, w=2.35, h=1.05, accent=BLUE):
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, LINE)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent)
    _textbox(slide, title, x + 0.18, y + 0.12, w - 0.25, 0.18, 8.2, GRAY, True)
    _textbox(slide, value, x + 0.18, y + 0.36, w - 0.25, 0.28, 18, NAVY, True)
    _textbox(slide, subtitle, x + 0.18, y + 0.76, w - 0.25, 0.16, 7.5, GRAY)


def _pill(slide, text, x, y, w, color):
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.32, color)
    _textbox(slide, text, x, y + 0.07, w, 0.12, 7.5, WHITE, True, "center")


def _bullets(slide, items, x, y, w, h, size=12, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0


def _table(slide, headers, rows, x, y, widths, row_h=0.42):
    total_w = sum(widths)
    total_h = row_h * (len(rows) + 1)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, total_w, total_h, WHITE, LINE)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, total_w, row_h, NAVY)

    cur = x
    for i, h in enumerate(headers):
        _textbox(slide, h, cur + 0.07, y + 0.12, widths[i] - 0.12, 0.16, 8, WHITE, True)
        cur += widths[i]

    for r, row in enumerate(rows):
        yy = y + row_h * (r + 1)
        if r % 2 == 1:
            _shape(slide, MSO_SHAPE.RECTANGLE, x, yy, total_w, row_h, RGBColor(241, 245, 249))
        cur = x
        for c, cell in enumerate(row):
            _textbox(slide, cell, cur + 0.07, yy + 0.11, widths[c] - 0.12, 0.14, 7.6, SLATE)
            cur += widths[c]


def _bar(slide, label, value, max_value, x, y, w, color):
    _textbox(slide, label, x, y - 0.02, 2.6, 0.15, 8, SLATE, True)
    _shape(slide, MSO_SHAPE.RECTANGLE, x + 2.65, y + 0.04, w, 0.14, RGBColor(226, 232, 240))
    ratio = max(0, min(1, _float(value) / _float(max_value))) if max_value else 0
    _shape(slide, MSO_SHAPE.RECTANGLE, x + 2.65, y + 0.04, w * ratio, 0.14, color)
    _textbox(slide, _pct(value), x + 2.65 + w + 0.15, y - 0.03, 0.8, 0.15, 8, NAVY, True)


def _roadmap_card(slide, period, title, body, x, color):
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.72, 3.75, 3.35, WHITE, LINE)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, 1.72, 3.75, 0.12, color)
    _textbox(slide, period, x + 0.25, 2.05, 2.8, 0.22, 18, color, True)
    _textbox(slide, title, x + 0.25, 2.55, 3.1, 0.28, 14, NAVY, True)
    _textbox(slide, body, x + 0.25, 3.05, 3.1, 1.2, 11, SLATE)
    _pill(slide, "MANAGEMENT GATE", x + 0.25, 4.55, 1.65, color)


def create_enterprise_pptx(
    company_name,
    sector,
    score,
    maturity,
    company_metrics,
    financial_result,
    recommendations,
    output_file="OptiFlow_Enterprise_Deck.pptx"
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    wait_rate = _safe(company_metrics, "wait_rate", 0)
    oee = _safe(company_metrics, "oee", 0)
    defect_rate = _safe(company_metrics, "defect_rate", 0)
    line_balance_loss = _safe(company_metrics, "line_balance_loss", 0)

    yearly_saving = _safe(financial_result, "Tahmini Yıllık Tasarruf", _safe(financial_result, "İyileştirme Potansiyeli", 0))
    total_loss = _safe(financial_result, "Toplam Operasyonel Kayıp", 0)
    roi = _safe(financial_result, "ROI (%)", 0)
    payback = _safe(financial_result, "Geri Dönüş Süresi (Ay)", 0)

    risk_score = round((_float(wait_rate) + _float(line_balance_loss) + _float(defect_rate) + max(0, 100 - _float(oee))) / 4, 1)
    risk_level = "Yüksek Risk" if risk_score >= 25 else "Orta Risk" if risk_score >= 15 else "Düşük Risk"

    page = 1

    # 1 COVER
    slide = _blank(prs)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, BLUE)
    _shape(slide, MSO_SHAPE.RECTANGLE, 9.4, 0, 3.95, 7.5, RGBColor(30, 64, 175))
    _shape(slide, MSO_SHAPE.RECTANGLE, 9.7, 0, 0.05, 7.5, CYAN)

    _textbox(slide, "OPTIFLOW", 0.85, 0.85, 4.5, 0.35, 24, WHITE, True)
    _textbox(slide, "Enterprise Operational Excellence Assessment", 0.87, 1.28, 6.8, 0.35, 15, RGBColor(203, 213, 225))
    _textbox(slide, company_name, 0.85, 2.55, 8.2, 0.75, 34, WHITE, True)
    _textbox(slide, f"{sector} | Executive Consulting Deck", 0.88, 3.32, 8.2, 0.28, 14, RGBColor(203, 213, 225))
    _textbox(slide, f"Prepared by OptiFlow Consulting | {date.today()}", 0.88, 6.65, 7.5, 0.18, 9, RGBColor(203, 213, 225))
    _pill(slide, "CONFIDENTIAL", 10.85, 0.85, 1.45, BLUE)

    # 2 EXECUTIVE SUMMARY
    page += 1
    slide = _blank(prs)
    _background(slide)
    _header(slide, "Executive Summary", "Financial impact, operational risk and improvement opportunity", page)

    _metric_card(slide, "OptiFlow Score", f"{score}/100", "Overall score", 0.65, 1.45, accent=BLUE)
    _metric_card(slide, "Risk Level", risk_level, f"{risk_score}/100", 3.15, 1.45, accent=RED if risk_score >= 25 else AMBER)
    _metric_card(slide, "Annual Saving", _money(yearly_saving), "Estimated benefit", 5.65, 1.45, accent=GREEN)
    _metric_card(slide, "ROI", _pct(roi), "Return potential", 8.15, 1.45, accent=_color_for("roi", roi))
    _metric_card(slide, "Payback", f"{payback} ay", "Recovery period", 10.65, 1.45, accent=BLUE)

    _textbox(slide, "Management Interpretation", 0.7, 3.05, 5.2, 0.25, 14, NAVY, True)
    _textbox(
        slide,
        f"{company_name} için en kritik iyileştirme alanı bekleme süreleri, hat dengeleme ve OEE yönetimidir. "
        f"Tahmini yıllık tasarruf potansiyeli {_money(yearly_saving)} seviyesindedir.",
        0.7,
        3.43,
        11.6,
        0.7,
        12.5,
        SLATE
    )

    _table(
        slide,
        ["Priority", "Management Focus", "Expected Outcome"],
        [
            ["1", "Bekleme azaltma", "Hızlı verimlilik kazanımı"],
            ["2", "Hat dengeleme", "Kapasite ve akış iyileşmesi"],
            ["3", "KPI ritmi", "Sürdürülebilir yönetim sistemi"]
        ],
        0.7,
        4.65,
        [1.2, 4.4, 5.2],
        0.48
    )

    # 3 KPI HEALTH
    page += 1
    slide = _blank(prs)
    _background(slide)
    _header(slide, "KPI Health Check", "Operational performance indicators and status view", page)

    kpis = [
        ("Bekleme Oranı", wait_rate, "wait", "Lower is better"),
        ("OEE", oee, "oee", "Higher is better"),
        ("Hata / Fire Oranı", defect_rate, "defect", "Lower is better"),
        ("Hat Denge Kaybı", line_balance_loss, "balance", "Lower is better"),
        ("ROI", roi, "roi", "Higher is better"),
    ]

    y = 1.55
    for label, value, key, note in kpis:
        color = _color_for(key, value)
        _bar(slide, label, value, 100, 0.75, y, 6.6, color)
        _pill(slide, _status_text(key, value), 10.15, y - 0.08, 1.3, color)
        _textbox(slide, note, 11.62, y - 0.01, 1.25, 0.14, 7.2, GRAY)
        y += 0.7

    _textbox(slide, "Executive note", 0.75, 5.65, 2.5, 0.22, 13, NAVY, True)
    _textbox(
        slide,
        "Kırmızı alanlar ilk 30 günlük aksiyon planında önceliklendirilmelidir. Sarı alanlar izleme ve kontrol gerektirir. Yeşil alanlar mevcut performansın korunması gereken güçlü alanlarıdır.",
        0.75,
        5.98,
        11.4,
        0.55,
        11.5,
        SLATE
    )

    # 4 FINANCIAL
    page += 1
    slide = _blank(prs)
    _background(slide)
    _header(slide, "Financial Impact & ROI", "Estimated loss, savings potential and payback view", page)

    _metric_card(slide, "Total Operational Loss", _money(total_loss), "Estimated annual loss", 0.75, 1.45, 2.9, 1.16, RED)
    _metric_card(slide, "Improvement Potential", _money(yearly_saving), "Annual savings estimate", 3.95, 1.45, 2.9, 1.16, GREEN)
    _metric_card(slide, "ROI", _pct(roi), "Return on improvement", 7.15, 1.45, 2.35, 1.16, _color_for("roi", roi))
    _metric_card(slide, "Payback", f"{payback} ay", "Estimated recovery", 9.85, 1.45, 2.35, 1.16, BLUE)

    _table(
        slide,
        ["Financial Driver", "Business Meaning", "Management Action"],
        [
            ["Bekleme kaybı", "Boşa geçen işçilik ve kapasite", "Akış ve darboğaz çalışması"],
            ["Fire / kalite kaybı", "Yeniden işleme ve müşteri riski", "Standart iş ve kalite kontrol"],
            ["Kapasite kaybı", "Ek yatırım ihtiyacı algısı", "Hat dengeleme ve OEE yönetimi"],
            ["ROI", "Yatırım kararını destekler", "Öncelikli proje portföyü"]
        ],
        0.75,
        3.12,
        [2.7, 4.5, 4.7],
        0.5
    )

    # 5 RISK
    page += 1
    slide = _blank(prs)
    _background(slide)
    _header(slide, "Risk & Root Cause View", "Delivery, quality, capacity and financial risk prioritization", page)

    _metric_card(slide, "Risk Score", f"{risk_score}/100", risk_level, 0.75, 1.45, 2.8, 1.15, RED if risk_score >= 25 else AMBER)
    _metric_card(slide, "Delivery Risk", _pct(wait_rate), "Waiting-driven risk", 3.9, 1.45, 2.8, 1.15, _color_for("wait", wait_rate))
    _metric_card(slide, "Quality Risk", _pct(defect_rate), "Defect-driven risk", 7.05, 1.45, 2.8, 1.15, _color_for("defect", defect_rate))
    _metric_card(slide, "Capacity Risk", _pct(line_balance_loss), "Balance-driven risk", 10.2, 1.45, 2.8, 1.15, _color_for("balance", line_balance_loss))

    _table(
        slide,
        ["5M Area", "Likely Root Cause", "Recommended Intervention"],
        [
            ["İnsan", "İş yükü dengesizliği", "Beceri matrisi ve vardiya performans takibi"],
            ["Makine", "Duruş ve kapasite kayıpları", "OEE bileşen takibi ve bakım planı"],
            ["Metot", "Standart iş eksikliği", "İş etüdü ve hat dengeleme"],
            ["Malzeme", "Akış gecikmeleri", "Kanban ve süreç içi stok kontrolü"],
            ["Ölçüm", "KPI görünürlüğü yetersiz", "Günlük KPI panosu ve yönetim ritmi"]
        ],
        0.75,
        3.25,
        [1.7, 4.7, 5.6],
        0.42
    )

    # 6 ROADMAP
    page += 1
    slide = _blank(prs)
    _background(slide)
    _header(slide, "30-60-90 Day Transformation Roadmap", "Commercially focused implementation plan", page)

    _roadmap_card(
        slide,
        "0-30 Gün",
        "Data Validation & Quick Wins",
        "Veri doğrulama, darboğaz analizi, KPI baz çizgisi ve hızlı kazanım alanlarının belirlenmesi.",
        0.75,
        BLUE
    )
    _roadmap_card(
        slide,
        "30-60 Gün",
        "Flow & Capacity Improvement",
        "Hat dengeleme, bekleme azaltma, standart iş ve OEE takip sisteminin devreye alınması.",
        4.8,
        AMBER
    )
    _roadmap_card(
        slide,
        "60-90 Gün",
        "Management System",
        "KPI toplantı ritmi, sorumluluk matrisi, performans panosu ve sürdürülebilir iyileştirme sistemi.",
        8.85,
        GREEN
    )

    # 7 ACTION PLAN
    page += 1
    slide = _blank(prs)
    _background(slide)
    _header(slide, "Management Action Plan", "Prioritized initiatives for executive decision making", page)

    action_rows = []
    default_actions = [
        "Bekleme sürelerini azalt",
        "Hat dengeleme çalışması başlat",
        "OEE takip sistemini kur",
        "KPI yönetim ritmi oluştur",
        "Kalite kayıplarının finansal etkisini izle"
    ]

    recs = recommendations if recommendations else default_actions

    for i, rec in enumerate(recs[:6], start=1):
        action_rows.append([str(i), str(rec), "High", "0-90 gün"])

    _table(
        slide,
        ["#", "Action", "Impact", "Timing"],
        action_rows,
        0.75,
        1.55,
        [0.6, 7.0, 1.8, 2.0],
        0.48
    )

    _textbox(slide, "Decision request", 0.75, 5.5, 3.0, 0.24, 13, NAVY, True)
    _textbox(
        slide,
        "Yönetimin ilk aşamada veri doğrulama, darboğaz analizi ve hat dengeleme çalışmalarını başlatması önerilir. Bu yaklaşım, yüksek maliyetli yatırım kararlarından önce mevcut sistemdeki kayıpları azaltmaya odaklanır.",
        0.75,
        5.86,
        11.5,
        0.55,
        11.5,
        SLATE
    )

    # 8 CLOSING
    page += 1
    slide = _blank(prs)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, BLUE)
    _textbox(slide, "Next Step", 0.85, 0.95, 6.0, 0.45, 30, WHITE, True)
    _textbox(
        slide,
        "Start a 30-day operational diagnostic sprint focused on data validation, bottleneck confirmation and quick-win implementation.",
        0.88,
        1.65,
        8.5,
        0.65,
        16,
        RGBColor(203, 213, 225)
    )

    _table(
        slide,
        ["Sprint", "Objective", "Output"],
        [
            ["Week 1", "Veri doğrulama", "KPI baseline"],
            ["Week 2", "Darboğaz gözlemi", "Critical loss map"],
            ["Week 3", "Hızlı kazanımlar", "Action portfolio"],
            ["Week 4", "Yönetim sunumu", "Implementation roadmap"]
        ],
        0.9,
        3.0,
        [1.6, 4.0, 4.7],
        0.45
    )
    _textbox(slide, "Prepared by OptiFlow Consulting", 0.9, 6.75, 4.0, 0.2, 9, RGBColor(203, 213, 225))

    prs.save(output_file)
    return output_file
